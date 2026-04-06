# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt
import os

W, H = 33.867, 19.05

def apply_style(run, size, bold=False, color=None):
    run.font.name = "微软雅黑"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)

def textbox(slide, l, t, w, h, text, size=12, bold=False, color=(0,0,0), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    shp = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    apply_style(r, size, bold, color)
    return shp

def rounded_box(slide, l, t, w, h, fill, border=(100,100,100), title=None, title_size=12, text_color=(0,0,0)):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*border)
    shp.line.width = Pt(1.0)
    if title:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        apply_style(r, title_size, True, text_color)
    return shp

def arrow(slide, x1, y1, x2, y2, color=(100,100,100), width=1.5, arrowhead=True):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    if arrowhead:
        line.line.end_arrowhead = True
    return line

def draw_lstm_chain(slide):
    # Modified to match the user's specific BiLSTM diagram structure
    xs = [11.8, 16.5, 21.2] # Center points for t-1, t, t+1
    y_in = 8.4
    y_fw = 7.4
    y_bw = 6.4
    y_out = 5.4
    y_out_label = 4.8
    
    colors = {"fw": (145, 199, 138), "bw": (224, 152, 160), "in": (255, 235, 205), "out": (255, 192, 203)}
    
    bw_h = 0.7
    bw_w = 1.8
    r_in = 0.5
    r_out = 0.4
    
    # Draw arrows FIRST so they are visually behind the boxes
    for x in xs:
        # x to fw
        arrow(slide, x, y_in-r_in, x, y_fw+bw_h/2)
        # x to bw (bypass fw on the right)
        l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x+0.3), Cm(y_in-r_in+0.1), Cm(x+0.9), Cm(y_fw))
        l1.line.color.rgb = RGBColor(100,100,100); l1.line.width = Pt(1.5)
        arrow(slide, x+0.9, y_fw, x+0.3, y_bw+bw_h/2)
        # fw to out (bypass bw on the left)
        l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x-0.3), Cm(y_fw-bw_h/2), Cm(x-0.9), Cm(y_bw))
        l2.line.color.rgb = RGBColor(100,100,100); l2.line.width = Pt(1.5)
        arrow(slide, x-0.9, y_bw, x-0.3, y_out+r_out)
        # bw to out
        arrow(slide, x, y_bw-bw_h/2, x, y_out+r_out)
        # out to y
        arrow(slide, x, y_out-r_out, x, y_out_label+0.4)
        
    # Horizontal arrows for FW (Left to Right)
    for i in range(2):
        arrow(slide, xs[i]+bw_w/2, y_fw, xs[i+1]-bw_w/2, y_fw, color=colors["fw"], width=2.0)
    # Horizontal arrows for BW (Right to Left)
    for i in range(2, 0, -1):
        arrow(slide, xs[i]-bw_w/2, y_bw, xs[i-1]+bw_w/2, y_bw, color=colors["bw"], width=2.0)

    # Input/Output text labels on left
    textbox(slide, 9.6, y_in-0.4, 1.8, 0.8, "输入层", size=10, bold=True)
    textbox(slide, 9.6, y_bw+0.1, 1.8, 0.8, "双向层\n(BiLSTM)", size=10, bold=True)
    textbox(slide, 9.6, y_out-0.4, 1.8, 0.8, "输出层", size=10, bold=True)
    
    # Draw Nodes
    labels = ["t-1", "t", "t+1"]
    for i, x in enumerate(xs):
        # Input circles
        c_in = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x-r_in), Cm(y_in-r_in), Cm(r_in*2), Cm(r_in*2))
        c_in.fill.solid(); c_in.fill.fore_color.rgb = RGBColor(*colors["in"])
        c_in.line.color.rgb = RGBColor(100,100,100)
        # Use wider textbox to avoid wrap
        textbox(slide, x-1.0, y_in-r_in, 2.0, r_in*2, f"x_{labels[i]}", size=10, bold=True)
        
        # Forward LSTM box
        fw_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x-bw_w/2), Cm(y_fw-bw_h/2), Cm(bw_w), Cm(bw_h))
        fw_box.fill.solid(); fw_box.fill.fore_color.rgb = RGBColor(220, 235, 255)
        fw_box.line.color.rgb = RGBColor(100,100,100)
        textbox(slide, x-bw_w/2, y_fw-bw_h/2, bw_w, bw_h, "LSTM", size=10, bold=True)
        
        # Backward LSTM box
        bw_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x-bw_w/2), Cm(y_bw-bw_h/2), Cm(bw_w), Cm(bw_h))
        bw_box.fill.solid(); bw_box.fill.fore_color.rgb = RGBColor(220, 235, 255)
        bw_box.line.color.rgb = RGBColor(100,100,100)
        textbox(slide, x-bw_w/2, y_bw-bw_h/2, bw_w, bw_h, "LSTM", size=10, bold=True)
        
        # Output + circle
        c_out = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x-r_out), Cm(y_out-r_out), Cm(r_out*2), Cm(r_out*2))
        c_out.fill.solid(); c_out.fill.fore_color.rgb = RGBColor(*colors["out"])
        c_out.line.color.rgb = RGBColor(100,100,100)
        textbox(slide, x-r_out, y_out-r_out, r_out*2, r_out*2, "+", size=14, bold=True)
        
        # Y labels (use wide textbox to avoid wrap)
        textbox(slide, x-1.0, y_out_label, 2.0, 0.6, f"y_{labels[i]}", size=10, bold=True)

    # Link to h_target
    arrow(slide, xs[-1]+bw_w/2+0.1, y_out, xs[-1]+bw_w/2+0.5, y_out)
    # Use wider box so h_target doesn't wrap awkwardly
    rounded_box(slide, xs[-1]+bw_w/2+0.5, y_out-0.4, 2.2, 0.8, (240, 240, 250), title="h_target\n(上下文)", title_size=9)

def main():
    prs = Presentation()
    prs.slide_width = Cm(W)
    prs.slide_height = Cm(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Colors
    c_bg = (248, 250, 253)
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(*c_bg)
    
    c_navy = (32, 58, 92)
    c_panel = (236, 242, 249)
    c_line = (116, 143, 176)
    
    # Title
    textbox(slide, 1, 0.5, 31, 1.2, "FA-BiLSTM 测井智能识别架构图", size=24, bold=True, color=c_navy)

    # Panels
    rounded_box(slide, 1, 2.5, 7.5, 15.5, c_panel, c_line)
    textbox(slide, 1, 2.5, 7.5, 1.2, "1. 数据输入 (Data Input)", size=16, bold=True, color=c_navy)
    
    rounded_box(slide, 9, 2.5, 15.5, 15.5, c_panel, c_line)
    textbox(slide, 9, 2.5, 15.5, 1.2, "2. 核心网络 (FA-BiLSTM Engine)", size=16, bold=True, color=c_navy)

    rounded_box(slide, 25, 2.5, 7.8, 15.5, c_panel, c_line)
    textbox(slide, 25, 2.5, 7.8, 1.2, "3. 预测输出 (Prediction Output)", size=16, bold=True, color=c_navy)
    
    # Arrows between panels
    arrow(slide, 8.5, 10, 9, 10, width=3.0)
    arrow(slide, 24.5, 10, 25, 10, width=3.0)

    # ------------------
    # 1. Left Panel
    # ------------------
    textbox(slide, 1.5, 4, 6.5, 1, "多维度连续测井参数", size=14, bold=True, color=(0,0,0))
    rounded_box(slide, 2.2, 5.5, 5, 1, (255,255,255), title="Sigma (俘获截面)", title_size=11)
    rounded_box(slide, 2.2, 7.0, 5, 1, (255,255,255), title="C/O (碳氧比)", title_size=11)
    rounded_box(slide, 2.2, 8.5, 5, 1, (255,255,255), title="Ratio (非弹性比)", title_size=11)
    
    # Slide Window Box
    rounded_box(slide, 1.5, 10.5, 6.5, 6, (230, 240, 255), title="")
    textbox(slide, 1.5, 10.8, 6.5, 1, "基于滑动窗口的时序构造", size=12, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.8, 12, 5.9, 4, "Window Size (L) = 10\n\n通过沿井深滑动，将单点转化为包含上下文的时序片段：\n(Batch, L, Features)", size=11, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP)

    # ------------------
    # 2. Middle Panel
    # ------------------
    # Top: BiLSTM
    rounded_box(slide, 9.5, 4, 14.5, 5.8, (255, 255, 255), c_line)
    textbox(slide, 9.5, 4.0, 14.5, 0.8, "时序特征提取 (BiLSTM Layer)", size=14, bold=True, color=c_navy)
    draw_lstm_chain(slide)
    textbox(slide, 9.5, 9.0, 14.5, 0.8, "说明：双向LSTM分别捕捉“由浅至深”和“由深至浅”的地质趋势", size=11, color=(100,100,100))

    # Bottom: Attention
    rounded_box(slide, 9.5, 10, 14.5, 7.5, (255, 255, 255), c_line)
    textbox(slide, 9.5, 10, 14.5, 1, "特征注意力门控 (Feature Attention Gate)", size=14, bold=True, color=c_navy)
    
    # Attention diagram
    rounded_box(slide, 10, 12, 2.8, 1.2, (217, 238, 243), title="h_target\n(时序上下文表示)", title_size=11)
    rounded_box(slide, 10, 15.5, 2.8, 1.2, (247, 231, 220), title="x_target\n(当前深度物理特征)", title_size=11)
    
    arrow(slide, 12.8, 12.6, 14.5, 13.8)
    arrow(slide, 12.8, 16.1, 14.5, 14.8)
    
    c_att = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(14.5), Cm(13.5), Cm(1.6), Cm(1.6))
    c_att.fill.solid(); c_att.fill.fore_color.rgb = RGBColor(233, 227, 241)
    c_att.line.color.rgb = RGBColor(*c_line)
    textbox(slide, 14.5, 13.5, 1.6, 1.6, "评分\n函数", size=10, bold=True, align=PP_ALIGN.CENTER)
    
    arrow(slide, 16.1, 14.3, 17, 14.3)
    
    rounded_box(slide, 17, 13.5, 2.5, 1.6, (228, 239, 224), title="Softmax激活\n分配权重(α)", title_size=11)
    
    arrow(slide, 19.5, 14.3, 20.8, 14.3)
    
    # Route x_target to output
    arrow(slide, 11.4, 16.7, 11.4, 17.1, arrowhead=False)
    arrow(slide, 11.4, 17.1, 21.8, 17.1, arrowhead=False)
    arrow(slide, 21.8, 17.1, 21.8, 15.1)
    
    rounded_box(slide, 20.8, 13.5, 2.5, 1.6, (255, 245, 230), title="特征贡献度加权\nα ⊙ x_target", title_size=11)
    
    textbox(slide, 10, 11.0, 13.5, 1, "逻辑：上下文 h_target 决定当前环境下，各个物理参数 x_target 的贡献大小", size=11, color=(100,100,100), align=PP_ALIGN.LEFT)

    # 上下文传递连线 (从 BiLSTM h_target 传递到 Attention h_target)
    # top h_target ( center x=23.7, y=5.4 ), bottom h_target ( center x=11.4, y=12.6 )
    conn1_1 = arrow(slide, 23.7, 5.8, 23.7, 9.8, color=(210, 150, 150), width=1.2, arrowhead=False)
    conn1_1.line.dash_style = 4
    conn1_2 = arrow(slide, 23.7, 9.8, 11.4, 9.8, color=(210, 150, 150), width=1.2, arrowhead=False)
    conn1_2.line.dash_style = 4
    conn1_3 = arrow(slide, 11.4, 9.8, 11.4, 12.0, color=(210, 150, 150), width=1.2, arrowhead=True)
    conn1_3.line.dash_style = 4
    textbox(slide, 21.0, 9.8, 2.5, 0.6, "上下文状态传递", size=10, bold=True, color=(180, 100, 100))

    # 共享特征连线 (从 BiLSTM输入 x_t 传递到 Attention x_target)
    # x_t center ~ x=16.5, y=8.4. Drop down to y=15.5 then go left to x=12.8
    # Bottom x_target right edge is at x=12.8, Box y is 15.5~16.7 center is 16.1
    conn2_1 = arrow(slide, 16.5, 8.9, 16.5, 16.1, color=(150, 180, 150), width=1.2, arrowhead=False)
    conn2_1.line.dash_style = 4
    conn2_2 = arrow(slide, 16.5, 16.1, 12.8, 16.1, color=(150, 180, 150), width=1.2, arrowhead=True)
    conn2_2.line.dash_style = 4
    textbox(slide, 15.3, 15.6, 2.5, 0.6, "当前时刻物理特征", size=10, bold=True, color=(100, 160, 100))


    # ------------------
    # 3. Right Panel
    # ------------------
    # Fusion
    rounded_box(slide, 25.5, 3.5, 6.8, 3.5, (255, 255, 255), c_line)
    textbox(slide, 25.5, 3.5, 6.8, 1, "多尺度特征融合", size=14, bold=True, color=c_navy)
    textbox(slide, 25.5, 4.5, 6.8, 2, "拼接上下文与加权特征：\n[ h_target,\n  x_target,\n  weighted_x ]\n送入全连接层", size=10, align=PP_ALIGN.CENTER)
    
    arrow(slide, 28.9, 7, 28.9, 8.5)

    # Output
    rounded_box(slide, 25.5, 8.5, 6.8, 9, (255, 255, 255), c_line)
    textbox(slide, 25.5, 8.5, 6.8, 1.2, "多阶段综合输出", size=14, bold=True, color=c_navy)
    
    # Output items
    rounded_box(slide, 26, 10.5, 5.8, 1.5, (250, 240, 240), title="1. 分类标签 (Gas/Oil/ 0/1)")
    rounded_box(slide, 26, 12.8, 5.8, 1.5, (240, 250, 240), title="2. 储层指数 (平滑概率曲线)")
    rounded_box(slide, 26, 15.1, 5.8, 1.5, (240, 240, 255), title="3. 参数敏感性 (α响应曲线)")

    # Save
    out_file = "FA_BiLSTM_单页架构图_一页版.pptx"
    prs.save(out_file)

if __name__ == "__main__":
    main()
