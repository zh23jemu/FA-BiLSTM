# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


W, H = 33.867, 19.05

COLORS = {
    "bg": (248, 250, 253),
    "navy": (32, 58, 92),
    "text": (43, 52, 64),
    "muted": (102, 113, 126),
    "line": (116, 143, 176),
    "panel": (236, 242, 249),
    "blue": (222, 234, 248),
    "green": (228, 239, 224),
    "peach": (247, 231, 220),
    "lavender": (233, 227, 241),
    "cyan": (217, 238, 243),
    "water": (149, 177, 220),
    "oil": (246, 171, 152),
    "gas": (177, 224, 235),
    "curve1": (104, 129, 162),
    "curve2": (94, 172, 182),
    "curve3": (221, 142, 120),
    "curve4": (166, 141, 199),
    "white": (255, 255, 255),
}


def style(run, size, bold=False, color=None):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*(color or COLORS["text"]))


def textbox(slide, l, t, w, h, text, size=12, bold=False, color=None,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style(r, size, bold, color)
    return shp


def rounded_box(slide, l, t, w, h, fill, title=None, subtitle=None, title_size=15, sub_size=10):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*COLORS["line"])
    shp.line.width = Pt(1.2)

    if title:
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title
        style(r1, title_size, True, COLORS["navy"])
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = subtitle
            style(r2, sub_size, False, COLORS["muted"])
    return shp


def arrow(slide, x1, y1, x2, y2, color=None, width=2, arrow_end=True):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    line.line.color.rgb = RGBColor(*(color or COLORS["line"]))
    line.line.width = Pt(width)
    if arrow_end:
        line.line.end_arrowhead = True
    return line


def panel(slide, l, t, w, h, title):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*COLORS["panel"])
    bg.line.color.rgb = RGBColor(*COLORS["line"])
    bg.line.width = Pt(1.1)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(0.85))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*COLORS["blue"])
    header.line.fill.background()
    textbox(slide, l, t + 0.08, w, 0.5, title, size=14, bold=True, color=COLORS["navy"],
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def mini_curve(slide, l, t, w, h, title, color):
    rounded_box(slide, l, t, w, h, (250, 252, 255))
    textbox(slide, l, t + 0.15, w, 0.35, title, size=10.5, bold=True, color=COLORS["text"], align=PP_ALIGN.CENTER)
    for i in range(1, 5):
        y = t + 0.65 + i * (h - 1.0) / 5
        arrow(slide, l + 0.12, y, l + w - 0.12, y, color=(224, 230, 237), width=0.8, arrow_end=False)
    for i in range(1, 4):
        x = l + 0.28 + i * (w - 0.56) / 4
        arrow(slide, x, t + 0.58, x, t + h - 0.2, color=(232, 237, 243), width=0.8, arrow_end=False)
    pts = [
        (l + 0.18, t + 0.9), (l + 0.42, t + 1.08), (l + 0.64, t + 0.95),
        (l + 0.88, t + 1.28), (l + 1.10, t + 1.15), (l + 1.35, t + 1.52),
        (l + 1.58, t + 1.35), (l + 1.82, t + 1.72), (l + 2.05, t + 1.56)
    ]
    for p1, p2 in zip(pts[:-1], pts[1:]):
        arrow(slide, p1[0], p1[1], p2[0], p2[1], color=color, width=2.0, arrow_end=False)
    textbox(slide, l - 0.04, t + h / 2 - 0.45, 0.28, 0.9, "深\n度", size=8.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def step_card(slide, l, t, w, h, fill, title, desc, step_no, kind):
    rounded_box(slide, l, t, w, h, fill)
    textbox(slide, l + 0.2, t + 0.15, w - 0.4, 0.45, f"步骤 {step_no}：{title}", size=12.5, bold=True,
            color=COLORS["navy"], align=PP_ALIGN.CENTER)

    # inner mini scene
    ix, iy, iw, ih = l + 0.5, t + 0.78, w - 1.0, 0.95
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(ix), Cm(iy), Cm(iw), Cm(ih))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(250, 252, 255)
    rect.line.color.rgb = RGBColor(*COLORS["line"])
    rect.line.width = Pt(0.8)

    if kind == "step1":
        for i, c in enumerate([COLORS["water"], COLORS["oil"], COLORS["gas"]]):
            band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(ix + 0.3), Cm(iy + 0.16 + i * 0.2), Cm(iw - 1.8), Cm(0.18))
            band.fill.solid(); band.fill.fore_color.rgb = RGBColor(*c); band.line.fill.background()
        target = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(ix + 2.2), Cm(iy + 0.17), Cm(0.9), Cm(0.46))
        target.fill.solid(); target.fill.fore_color.rgb = RGBColor(255, 255, 255); target.line.color.rgb = RGBColor(*COLORS["line"]); target.line.width = Pt(0.7)
        textbox(slide, ix + 2.23, iy + 0.28, 0.84, 0.16, "目标层段", size=8.5, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
        textbox(slide, ix + iw - 1.05, iy + 0.15, 0.85, 0.48, "h_\ntarget", size=9.5, bold=True, color=COLORS["text"], align=PP_ALIGN.CENTER)

    elif kind == "step2":
        upper = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(ix + 0.3), Cm(iy + 0.18), Cm(iw - 0.6), Cm(0.24))
        upper.fill.solid(); upper.fill.fore_color.rgb = RGBColor(*COLORS["gas"]); upper.line.fill.background()
        lower = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(ix + 0.3), Cm(iy + 0.42), Cm(iw - 0.6), Cm(0.24))
        lower.fill.solid(); lower.fill.fore_color.rgb = RGBColor(*COLORS["oil"]); lower.line.fill.background()
        alpha = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(ix + 2.35), Cm(iy + 0.1), Cm(0.72), Cm(0.72))
        alpha.fill.solid(); alpha.fill.fore_color.rgb = RGBColor(245, 247, 251); alpha.line.color.rgb = RGBColor(*COLORS["line"])
        textbox(slide, ix + 2.44, iy + 0.3, 0.54, 0.18, "α", size=11, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
        textbox(slide, ix + 1.0, iy + 0.2, 0.9, 0.18, "高权重", size=8.5, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
        textbox(slide, ix + 3.5, iy + 0.45, 0.9, 0.18, "低权重", size=8.5, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)

    elif kind == "step3":
        for x, txt, c in [(ix + 0.55, "分类", COLORS["gas"]), (ix + 1.95, "曲线", COLORS["oil"]), (ix + 3.35, "权重", COLORS["water"])]:
            r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(iy + 0.2), Cm(0.9), Cm(0.45))
            r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*c); r.line.fill.background()
            textbox(slide, x, iy + 0.28, 0.9, 0.16, txt, size=8.8, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)

    textbox(slide, l + 0.3, t + 1.78, w - 0.6, 0.38, desc, size=9.8, color=COLORS["muted"], align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Cm(W)
prs.slide_height = Cm(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(*COLORS["bg"])

# Title
textbox(slide, 2.0, 0.35, 30.0, 0.9, "基于 BiLSTM 与特征注意力机制的测井智能识别框架",
        size=22, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)

# Panels
panel(slide, 0.8, 2.0, 8.8, 13.8, "输入")
panel(slide, 10.2, 2.0, 12.6, 13.8, "核心网络")
panel(slide, 23.4, 2.0, 9.6, 13.8, "输出")

# Left panel
textbox(slide, 1.2, 3.0, 8.0, 0.55, "多参数测井输入", size=16, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
mini_curve(slide, 1.2, 4.0, 1.65, 7.0, "MLR", COLORS["curve1"])
mini_curve(slide, 3.05, 4.0, 1.65, 7.0, "AMPST", COLORS["curve2"])
mini_curve(slide, 4.90, 4.0, 1.65, 7.0, "PHIE", COLORS["curve3"])
mini_curve(slide, 6.75, 4.0, 1.65, 7.0, "RICX", COLORS["curve4"])
textbox(slide, 1.2, 11.35, 8.0, 1.95,
        "输入可以是 MLR、AMPST、PHIE 等储层参数，\n也可以是 Sigma、RICX、RIN13、RATO13 等流体参数。\n通过滑动窗口构造长度为 10 的时序样本。",
        size=10.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)

# Middle top: BiLSTM
rounded_box(slide, 11.0, 3.05, 11.0, 4.35, COLORS["blue"])
textbox(slide, 11.3, 3.35, 10.4, 0.45, "BiLSTM 时序特征提取", size=16, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
textbox(slide, 11.6, 4.15, 9.8, 0.35, "输出末端时刻上下文表示  h_target", size=12, bold=True, color=COLORS["text"], align=PP_ALIGN.CENTER)
xs = [12.15, 13.65, 15.15, 16.65, 18.15]
for x in xs:
    shp1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x), Cm(4.95), Cm(0.72), Cm(0.72))
    shp1.fill.solid(); shp1.fill.fore_color.rgb = RGBColor(145, 199, 138); shp1.line.color.rgb = RGBColor(255, 255, 255)
    shp2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x), Cm(6.1), Cm(0.72), Cm(0.72))
    shp2.fill.solid(); shp2.fill.fore_color.rgb = RGBColor(122, 170, 214); shp2.line.color.rgb = RGBColor(255, 255, 255)
for i in range(len(xs) - 1):
    arrow(slide, xs[i] + 0.72, 5.3, xs[i + 1], 5.3, color=(145, 199, 138), width=1.5)
    arrow(slide, xs[i + 1], 6.45, xs[i] + 0.72, 6.45, color=(122, 170, 214), width=1.5)
textbox(slide, 11.8, 5.95, 9.4, 0.45, "1层双向LSTM  |  hidden_size = 16  |  输出维度 = 32", size=10.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)

# Middle bottom: attention
rounded_box(slide, 11.0, 8.3, 11.0, 5.7, COLORS["green"])
textbox(slide, 11.25, 8.6, 10.5, 0.45, "特征注意力耦合机制", size=16, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
rounded_box(slide, 12.0, 10.55, 2.15, 1.65, COLORS["peach"], "当前特征", "x_target", title_size=12, sub_size=10)
rounded_box(slide, 15.0, 10.25, 2.65, 2.25, COLORS["white"], "加权输出", "weighted_x\n= α ⊙\nx_target", title_size=12, sub_size=10)
rounded_box(slide, 18.95, 9.35, 1.95, 1.55, COLORS["cyan"], "上下文", "h_target", title_size=11.5, sub_size=10)
rounded_box(slide, 18.95, 12.1, 1.95, 1.55, COLORS["lavender"], "特征权重", "α", title_size=11.5, sub_size=12)
arrow(slide, 14.15, 11.35, 15.0, 11.35, color=(214, 131, 118), width=1.8)
arrow(slide, 18.95, 10.1, 17.65, 10.9, color=(103, 177, 191), width=1.8)
arrow(slide, 18.95, 12.85, 17.65, 11.95, color=(121, 140, 205), width=1.8)
textbox(slide, 11.65, 13.25, 9.7, 0.75,
        "注意力机制由 h_target 与 x_target 共同生成特征权重 α，\n并通过加权形成 weighted_x，实现多参数耦合表达。",
        size=10.2, color=COLORS["muted"], align=PP_ALIGN.CENTER)

# Right steps
step_card(slide, 24.05, 3.45, 8.0, 2.35, COLORS["peach"], "时序判别", "BiLSTM 提取上下文时序特征", 1, "step1")
step_card(slide, 24.05, 6.55, 8.0, 2.35, COLORS["green"], "特征加权", "注意力机制计算 α 并生成 weighted_x", 2, "step2")
step_card(slide, 24.05, 9.65, 8.0, 2.35, COLORS["lavender"], "融合输出", "输出分类结果、连续曲线与权重解释", 3, "step3")
arrow(slide, 28.05, 5.8, 28.05, 6.55, width=1.7)
arrow(slide, 28.05, 8.9, 28.05, 9.65, width=1.7)
textbox(slide, 24.3, 12.5, 7.5, 0.8, "典型输出包括：有效储层识别、气液识别、FA_alpha_* 权重曲线",
        size=10.2, color=COLORS["muted"], align=PP_ALIGN.CENTER)

# Inter-panel arrows
arrow(slide, 8.45, 12.4, 11.0, 12.4, color=COLORS["curve1"], width=2.0)
arrow(slide, 22.0, 5.8, 24.05, 4.65, width=2.0)
arrow(slide, 22.0, 11.2, 24.05, 7.7, width=2.0)

# Output bar
bar_x = 32.35
for y, hh, c, lab in [(3.95, 2.15, COLORS["gas"], "气"), (6.1, 2.1, COLORS["oil"], "储"), (8.2, 2.45, COLORS["water"], "层")]:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(bar_x), Cm(y), Cm(0.85), Cm(hh))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*c)
    rect.line.color.rgb = RGBColor(*COLORS["line"])
    rect.line.width = Pt(0.8)
    textbox(slide, bar_x + 0.05, y + hh / 2 - 0.18, 0.75, 0.36, lab, size=11, bold=True,
            color=COLORS["navy"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
textbox(slide, 31.95, 10.95, 1.85, 0.9, "结果柱状\n示意", size=9.2, color=COLORS["muted"], align=PP_ALIGN.CENTER)

# Footer
textbox(slide, 2.0, 16.55, 30.0, 0.8,
        "说明：BiLSTM 负责建模井深方向上的时序依赖关系，特征注意力层负责计算多参数动态权重，并通过加权与融合实现多参数耦合表达。",
        size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)

out = "FA_BiLSTM_网络结构图_论文风版.pptx"
prs.save(out)
print(out)
