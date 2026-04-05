# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt
import os

# Standard 16:9 Slide Size
W, H = 33.867, 19.05

COLORS = {
    "bg": (255, 255, 255),
    "navy": (26, 54, 104),
    "header_bg": (240, 244, 250),
    "text": (33, 37, 41),
    "muted": (108, 117, 125),
    "line": (73, 80, 87),
    "blue": (0, 123, 255),
    "light_blue": (230, 240, 255),
    "green": (40, 167, 69),
    "light_green": (230, 255, 230),
    "orange": (253, 126, 20),
    "light_orange": (255, 240, 230),
    "purple": (111, 66, 193),
    "light_purple": (240, 230, 255),
    "white": (255, 255, 255),
}

def apply_style(run, size, bold=False, color=None):
    run.font.name = "微软雅黑"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*(color or COLORS["text"]))

def add_title_box(slide, text):
    # Header bar
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(W), Cm(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*COLORS["header_bg"])
    rect.line.fill.background()
    
    # Text
    shp = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(W-2), Cm(1.2))
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    apply_style(r, 24, bold=True, color=COLORS["navy"])

def add_content_textbox(slide, l, t, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    shp = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    apply_style(r, size, bold, color)
    return shp

def add_rounded_rect(slide, l, t, w, h, fill_color, border_color=None, text=None, text_size=12, text_bold=False):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill_color)
    shp.line.color.rgb = RGBColor(*(border_color or COLORS["line"]))
    shp.line.width = Pt(1)
    
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        apply_style(r, text_size, text_bold)
    return shp

def add_arrow(slide, x1, y1, x2, y2, color=None, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    line.line.color.rgb = RGBColor(*(color or COLORS["line"]))
    line.line.width = Pt(width)
    line.line.end_arrowhead = True

def create_ppt():
    prs = Presentation()
    prs.slide_width = Cm(W)
    prs.slide_height = Cm(H)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rounded_rect(slide, 4, 6, 25.867, 7, COLORS["header_bg"], border_color=COLORS["navy"])
    add_content_textbox(slide, 5, 7.5, 24, 2, "基于 FA-BiLSTM 的三阶段分层流体识别框架", size=32, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
    add_content_textbox(slide, 5, 10, 24, 1.5, "网络架构、注意力机制与数据流分析", size=18, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # 2. Overall Framework
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "1. 整体框架设计 (Overall Framework)")
    add_content_textbox(slide, 1, 2.5, 31, 1, "参考论文的三阶段识别逻辑，结合深度卷积时序特征与特征动态权重", size=14, color=COLORS["muted"])
    
    # Three panels
    add_rounded_rect(slide, 1, 4, 8, 13, COLORS["light_blue"], text="数据输入\n(测井曲线组合)", text_size=18, text_bold=True)
    add_rounded_rect(slide, 10, 4, 14, 13, COLORS["light_green"], text="核心引擎\n(FA-BiLSTM 网络)", text_size=18, text_bold=True)
    add_rounded_rect(slide, 25, 4, 8, 13, COLORS["light_orange"], text="三步识别流程\n(分级输出)", text_size=18, text_bold=True)
    
    add_arrow(slide, 9, 10.5, 10, 10.5, width=3)
    add_arrow(slide, 24, 10.5, 25, 10.5, width=3)

    # 3. Data Input
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "2. 输入：多维度测井参数与滑动窗口")
    add_content_textbox(slide, 2, 3, 29, 2, "程序输入包含多维度测井特征，通过滑动窗口捕捉井段垂直方向上的时序关联性。", size=16)
    
    # Features List
    add_rounded_rect(slide, 2, 6, 6, 4, COLORS["header_bg"], text="特征 A: Sigma\n(俘获截面)")
    add_rounded_rect(slide, 9, 6, 6, 4, COLORS["header_bg"], text="特征 B: C/O\n(碳氧比)")
    add_rounded_rect(slide, 16, 6, 6, 4, COLORS["header_bg"], text="特征 C: Ratio\n(非弹性/俘获比)")
    add_rounded_rect(slide, 23, 6, 8, 4, COLORS["header_bg"], text="... 等 10+ 维度\n(根据模型动态加载)")
    
    # Sliding Window
    add_rounded_rect(slide, 2, 11, 29, 6, COLORS["white"], border_color=COLORS["blue"])
    add_content_textbox(slide, 3, 11.5, 27, 4, "时序滑动窗口 (Sliding Window)\n\n• 窗口大小 (Window Size): 10\n• 作用：通过建模当前深度点及其上方 9 个采样点的组合特征，识别沉积环境与地质接触面的变化。\n• 数据形状: (Batch, 10, Features_Dim)", size=16)

    # 4. BiLSTM Layer
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "3. 主体：BiLSTM 时序特征提取网络")
    add_content_textbox(slide, 2, 2.8, 29, 2, "BiLSTM (Bidirectional LSTM) 是网络的核心部分，负责提取井深序列中的上下文特征。", size=16)
    
    # BiLSTM Diagram (simplified)
    for i in range(5):
        x = 4 + i * 5
        add_rounded_rect(slide, x, 6, 1.5, 1.5, COLORS["white"], text=f"X_{i}", text_size=10)
        add_rounded_rect(slide, x, 9, 2, 4, COLORS["blue"], text="BiLSTM cell", text_size=10)
        add_arrow(slide, x + 0.75, 7.5, x + 0.75, 9, width=1)
        if i < 4:
            add_arrow(slide, x + 2, 10.5, x + 5, 10.5, color=COLORS["blue"], width=2)
            add_arrow(slide, x + 5, 11.5, x + 2, 11.5, color=COLORS["purple"], width=2)

    add_content_textbox(slide, 2, 14, 29, 4, "代码实现 (PyTorch):\nself.lstm = nn.LSTM(input_size, hidden_size, bidirectional=True, batch_first=True)\n\n• bidirectional=True: 同时学习“由浅及深”和“由深及浅”的地层趋势信息。\n• hidden_size=16: 控制模型复杂度，平衡泛化能力与参数量。", size=14, color=COLORS["navy"])

    # 5. Attention Mechanism
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "4. 特征注意力机制 (Feature Attention Gate)")
    add_content_textbox(slide, 2, 2.5, 29, 2, "注意力机制不再对时域加权，而是对“特征维度”进行动态加权，识别当前层段最敏感的特征。", size=16)
    
    # Logic blocks
    add_rounded_rect(slide, 2, 5, 6, 3, COLORS["light_blue"], text="h_context\n(BiLSTM 时序特征)", text_size=14)
    add_rounded_rect(slide, 10, 5, 6, 3, COLORS["light_orange"], text="x_current\n(当前输入特征)", text_size=14)
    
    add_rounded_rect(slide, 2, 10, 14, 4, COLORS["light_purple"], text="计算评分 Score:\nv(tanh(W_h * h + W_x * x))", text_size=16, text_bold=True)
    add_arrow(slide, 5, 8, 5, 10)
    add_arrow(slide, 13, 8, 13, 10)
    
    add_arrow(slide, 16, 12, 18, 12)
    add_rounded_rect(slide, 18, 10, 12, 4, COLORS["light_green"], text="特征权重 α (Sigmoid/Softmax)\nα1=Sigma, α2=CO, α3=Ratio...", text_size=16)
    
    add_content_textbox(slide, 2, 15, 29, 3, "注意力作用机制：\nWeighted_X = α ⊙ x_current\n实现“特征贡献度”的自动排序。例如：在低孔隙度地层，模型自动提高 C/O 的权重 α2。")

    # 6. Output and Fusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "5. 输出：多目标识别结果")
    add_content_textbox(slide, 2, 2.5, 29, 4, "模型输出层将 BiLSTM 时序特征与注意力加权后的特征融合，得出最终结论。", size=16)
    
    add_rounded_rect(slide, 2, 5, 29, 3, COLORS["header_bg"], text="特征融合: [h_target, x_target, weighted_x]", text_size=18, text_bold=True)
    add_arrow(slide, 16.5, 8, 16.5, 10)
    
    # Outputs boxes
    add_rounded_rect(slide, 2, 10, 9, 5, COLORS["white"], text="1. 分类标签\n(Gas / Oil / Water)")
    add_rounded_rect(slide, 12, 10, 9, 5, COLORS["white"], text="2. 连续流体指数\n(I_RQI, I_GRI, I_OII)")
    add_rounded_rect(slide, 22, 10, 9, 5, COLORS["white"], text="3. 特征重要度 α\n(可解释性分析)")
    
    add_content_textbox(slide, 2, 16, 29, 2, "模型输出 shape: (Batch, Num_Classes) 和 (Batch, Num_Features)", size=14, color=COLORS["muted"])

    # 7. Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, "结论与亮点")
    add_content_textbox(slide, 3, 5, 27, 10, "• 核心架构: BiLSTM 时序建模 + Feature Attention 特征门控。\n• 输入: 多维度测井深度序列数据 (Window=10)。\n• 特征注意力: 解决了多参数测井响应在不同地层环境下的权重不确定性。\n• 输出: 涵盖分类、连续曲线与权重解释，提供全方位决策支持。", size=18)

    save_path = "BiLSTM_Attention_Network_Intro.pptx"
    prs.save(save_path)
    return save_path

if __name__ == "__main__":
    path = create_ppt()
    print(f"PPT generated at: {os.path.abspath(path)}")
